"""Text extraction → silver-layer :class:`RawDocument`.

Supports the three document families the brief calls for:

  * **text / markdown** — decoded directly (pure-Python, no deps).
  * **PDF** — via ``pypdf`` (pure-Python, fits a Lambda zip).
  * **PPTX** — via ``python-pptx`` (pulls native ``lxml``; ship in a Lambda
    container/layer, not a bare zip).

``pypdf`` and ``python-pptx`` are optional (the ``extract`` extra) and imported
lazily, so the core API image — which doesn't extract — stays slim. The single
entry point :func:`extract_document` dispatches on content type (falling back
to the filename extension, since S3 events often report
``application/octet-stream``), so every caller (the ingestion Lambda, the
``POST /ingest`` handler, the seed script) goes through one code path.
"""

from __future__ import annotations

import io

from fossilrag.logging import get_logger
from fossilrag.models import RawDocument

log = get_logger("ingest")

# Canonical content type → short extractor tag (recorded in provenance).
SUPPORTED_CONTENT_TYPES: dict[str, str] = {
    "text/plain": "txt",
    "text/markdown": "md",
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}

# Filename extension → canonical content type, for when the supplied content
# type is missing or a generic octet-stream.
_EXT_TO_CTYPE: dict[str, str] = {
    ".txt": "text/plain",
    ".text": "text/plain",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".pdf": "application/pdf",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# Content types we never trust to name the format (S3's default, generic).
_GENERIC = {"", "application/octet-stream", "binary/octet-stream"}


def _normalise_newlines(text: str) -> str:
    """CRLF/CR → LF so paragraph splitting works regardless of source OS."""
    return text.replace("\r\n", "\n").replace("\r", "\n")


def _resolve_content_type(filename: str, content_type: str) -> str:
    """Pick the canonical content type from the header, else the extension."""
    ctype = content_type.split(";", 1)[0].strip().lower()
    if ctype in SUPPORTED_CONTENT_TYPES:
        return ctype
    if ctype in _GENERIC:
        # Generic/blank header — infer from the extension.
        ext = _extension(filename)
        if ext in _EXT_TO_CTYPE:
            return _EXT_TO_CTYPE[ext]
    raise ValueError(
        f"Unsupported content_type {content_type!r} for {filename!r}. "
        f"Supported: {sorted(SUPPORTED_CONTENT_TYPES)}."
    )


def _extension(filename: str) -> str:
    dot = filename.rfind(".")
    return filename[dot:].lower() if dot != -1 else ""


def _extract_text_bytes(data: bytes) -> str:
    # utf-8-sig strips a BOM if present; errors="replace" keeps ingestion
    # robust to the occasional mis-encoded byte rather than failing a document.
    return _normalise_newlines(data.decode("utf-8-sig", errors="replace"))


def _extract_pdf(data: bytes) -> str:
    """Extract text from a PDF, one page per block, in page order."""
    from pypdf import PdfReader  # lazy: optional 'extract' dependency

    reader = PdfReader(io.BytesIO(data))
    pages = [(page.extract_text() or "").strip() for page in reader.pages]
    return _normalise_newlines("\n\n".join(p for p in pages if p))


def _extract_pptx(data: bytes) -> str:
    """Extract text from a PPTX: every shape's text frame, slide by slide.

    PR1 covers text-frame shapes (titles, bodies, text boxes). Table cells,
    grouped shapes, and chart text are deepened with the Slide Mutator (PR8),
    which needs structured slide content.
    """
    from pptx import Presentation  # lazy: optional 'extract' dependency

    prs = Presentation(io.BytesIO(data))
    blocks: list[str] = []
    for slide in prs.slides:
        parts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if parts:
            blocks.append("\n".join(parts))
    return _normalise_newlines("\n\n".join(blocks))


_EXTRACTORS = {
    "txt": _extract_text_bytes,
    "md": _extract_text_bytes,
    "pdf": _extract_pdf,
    "pptx": _extract_pptx,
}


def extract_document(
    *,
    filename: str,
    data: bytes,
    content_type: str = "text/plain",
    user_id: str | None = None,
    source_uri: str | None = None,
) -> RawDocument:
    """Extract text + metadata from raw bytes into a :class:`RawDocument`.

    Raises ``ValueError`` for unsupported types so callers return a clean 4xx
    rather than indexing garbage.
    """
    ctype = _resolve_content_type(filename, content_type)
    tag = SUPPORTED_CONTENT_TYPES[ctype]
    try:
        text = _EXTRACTORS[tag](data)
    except Exception as exc:  # noqa: BLE001 — normalise parse failures to ValueError
        # A *supported* type with a corrupt/mismatched body (truncated PDF,
        # non-zip .pptx, ...) raises library-specific errors (pypdf's
        # PdfReadError, python-pptx's PackageNotFoundError) that don't subclass
        # ValueError. Normalise so callers map it to a clean 4xx and the
        # ingestion Lambda treats it as a non-retryable bad object — not an
        # opaque 500 / endless retry.
        raise ValueError(f"failed to extract {ctype} from {filename!r}: {exc}") from exc

    doc = RawDocument.from_text(
        filename=filename,
        text=text,
        content_type=ctype,
        user_id=user_id,
        source_uri=source_uri,
        metadata={"extractor": tag, "bytes": str(len(data))},
    )
    if doc.char_count == 0:
        # e.g. a scanned/image-only PDF — surface it so empty ingests aren't
        # silently invisible downstream (chunking would produce nothing).
        log.warning(
            "event=document_empty doc_id=%s filename=%s extractor=%s",
            doc.doc_id[:12],
            filename,
            tag,
        )
    log.info(
        "event=document_extracted doc_id=%s filename=%s content_type=%s extractor=%s chars=%d",
        doc.doc_id[:12],
        filename,
        ctype,
        tag,
        doc.char_count,
    )
    return doc
